# -*- coding: utf-8 -*-
"""The PPD final presentation — a fresh deck in the thesis-defence template language.

NOT part of the V6a build. This deck is presented, not archived: 19 physical pages, low text
density, one idea per page, in the palette of the reference template he supplied
("Gray and White Simple Thesis Defense", Canva):

    cream F6F4F1 · ink 253439 · tan B29E84 · sage 7C898B

with its motifs — the left rail (dark block over a long sage bar), the tri-square corner mark,
the tricolour underline, vertical edge text — and the JU logo lifted from the Weekly deck's own
master. Lato is not installed on this machine, so Segoe UI Light / Segoe UI stand in.

Slide 5 is TWO pages on purpose: python-pptx cannot author animations, so "all 19 visible, then
three highlighted and the rest greyed on click" is done as consecutive slides — the click IS the
transition.

The three demo videos are re-encoded timelapses (15 s each) of real pulls, embedded in the file;
they regenerate from Test data if the cached MP4s are missing. Numbers on the validation pages
are the ones mot2_compare / mot2_pp print today, not remembered ones.

Build: python documentation/scripts/ppd_final_deck.py
Writes to the scratchpad first (OneDrive corrupts freshly written media parts), then the caller
copies into documentation/decks/.
"""
import os
import subprocess
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.abspath(os.path.join(HERE, "..", ".."))
FIGS = os.path.join(ROOT, "documentation", "figures")
SP = os.environ.get("PPD_SCRATCH", os.path.join(HERE, "_ppd_build"))
MEDIA = os.path.join(SP, "media")
OUT = os.path.join(SP, "PPD_Final_Presentation.pptx")

# ---- palette (measured from the reference pptx) and type
CREAM = RGBColor(0xF6, 0xF4, 0xF1)
INK = RGBColor(0x25, 0x34, 0x39)
TAN = RGBColor(0xB2, 0x9E, 0x84)
SAGE = RGBColor(0x7C, 0x89, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FAINT = RGBColor(0xDE, 0xDA, 0xD2)          # greyed-out tiles, hairlines on cream
DISP, BODY, SEMI = "Segoe UI Light", "Segoe UI", "Segoe UI Semibold"

EW, EH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(EW)
prs.slide_height = Inches(EH)
BLANK = prs.slide_layouts[6]


# ================================================================ drawing vocabulary
def rect(s, x, y, w, h, fill, line=None, lw=1.0):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line; r.line.width = Pt(lw)
    r.shadow.inherit = False
    return r


def text(s, x, y, w, h, runs, *, fs=12, font=BODY, colour=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, spacing=1.0, wrap=True):
    """One box, many paragraphs. `runs` is a string, or a list of paragraphs, each a string or a
    list of (txt, {font/fs/colour/bold/italic}) run tuples."""
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = [runs] if isinstance(runs, str) else runs
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        for run in ([(para, {})] if isinstance(para, str) else para):
            t, kw = run
            r = p.add_run(); r.text = t
            r.font.name = kw.get("font", font)
            r.font.size = Pt(kw.get("fs", fs))
            r.font.bold = kw.get("bold", bold)
            r.font.italic = kw.get("italic", False)
            r.font.color.rgb = kw.get("colour", colour)
    return box


def bullets(s, x, y, w, h, items, *, fs=11.5, gap=6, colour=INK):
    """Tan square markers, ink text. An item may be (text, dict-overrides)."""
    paras = []
    for it in items:
        t, kw = it if isinstance(it, tuple) else (it, {})
        paras.append([("▪  ", {"colour": TAN, "fs": kw.get("fs", fs), "bold": True}),
                      (t, {"fs": kw.get("fs", fs), "colour": kw.get("colour", colour),
                           "bold": kw.get("bold", False), "font": kw.get("font", BODY)})])
    box = text(s, x, y, w, h, paras, fs=fs, spacing=1.04)
    for p in box.text_frame.paragraphs:
        p.space_after = Pt(gap)
    return box


def card(s, x, y, w, h, *, fill=WHITE, line=FAINT, lw=1.0):
    return rect(s, x, y, w, h, fill, line=line, lw=lw)


def pic(s, path, x, y, maxw, maxh, *, frame=True, pad=0.10):
    """Image fitted and centred inside a white card."""
    iw, ih = Image.open(path).size
    availw, availh = maxw - 2 * pad, maxh - 2 * pad
    sc = min(availw / iw, availh / ih)
    w, h = iw * sc, ih * sc
    if frame:
        card(s, x, y, maxw, maxh)
    s.shapes.add_picture(path, Inches(x + (maxw - w) / 2), Inches(y + (maxh - h) / 2),
                         width=Inches(w))
    return h


def chip(s, x, y, txt, *, fill=INK, fg=WHITE, fs=9.5, w=None, h=0.30):
    w = w or (0.16 + 0.092 * len(txt))
    c = rect(s, x, y, w, h, fill)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.name = SEMI; r.font.size = Pt(fs); r.font.color.rgb = fg
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return c


# ---- the template furniture
def base(s, *, rail=True):
    rect(s, 0, 0, EW, EH, CREAM)
    if rail:
        rect(s, 0, 0, 0.16, 1.02, INK)
        rect(s, 0, 1.14, 0.16, EH - 1.14, SAGE)


def ju(s, x=0.46, y=0.28, hgt=0.44):
    # the real JU lockup (mark + wordmark), lifted from his own deck's layout media and
    # tinted to this template's ink - it carries its own transparency, nothing to compose
    s.shapes.add_picture(os.path.join(FIGS, "ju_mark.png"), Inches(x), Inches(y),
                         height=Inches(hgt))


def squares(s, x=12.32, y=0.38, a=0.13, gap=0.115):
    for i, c in enumerate((INK, TAN, SAGE)):
        rect(s, x + i * (a + gap), y, a, a, c)


def vside(s, txt):
    box = text(s, 12.70, 2.6, 3.0, 0.34, txt, fs=8.5, font=SEMI, colour=SAGE,
               align=PP_ALIGN.CENTER, wrap=False)
    box.rotation = 270
    box.left, box.top = Inches(11.62), Inches(3.62)   # rotate about centre, then place on the edge


def tricolour(s, x, y, w=3.6, h=0.055):
    rect(s, x, y, w * 0.26, h, INK)
    rect(s, x + w * 0.26, y, w * 0.13, h, TAN)
    rect(s, x + w * 0.39, y, w * 0.61, h, SAGE)


_PAGE = [0]


def pageno(s):
    _PAGE[0] += 1
    text(s, 12.55, 7.08, 0.62, 0.3, "%02d" % _PAGE[0], fs=10, font=SEMI, colour=SAGE,
         align=PP_ALIGN.RIGHT)


def header(s, kicker, title_txt, *, tfs=29):
    text(s, 0.85, 0.86, 11.4, 0.30, kicker.upper(), fs=10.5, font=SEMI, colour=TAN)
    text(s, 0.83, 1.10, 11.9, 0.62, title_txt, fs=tfs, font=DISP, colour=INK)
    tricolour(s, 0.85, 1.74)


def content_slide(kicker, title_txt, side=None, *, tfs=29):
    s = prs.slides.add_slide(BLANK)
    base(s)
    ju(s)
    squares(s)
    header(s, kicker, title_txt, tfs=tfs)
    vside(s, side or "PPD-UTM  ·  FINAL PRESENTATION")
    pageno(s)
    return s


def ghost(s, x, y, w, h, label, sub):
    """A dashed placeholder for content that arrives later - visibly intentional."""
    g = rect(s, x, y, w, h, CREAM, line=SAGE, lw=1.2)
    g.line.dash_style = 4                                    # dashed
    text(s, x, y + h / 2 - 0.34, w, 0.4, label, fs=13, font=SEMI, colour=SAGE,
         align=PP_ALIGN.CENTER)
    text(s, x + 0.3, y + h / 2 + 0.02, w - 0.6, 0.5, sub, fs=10, colour=SAGE,
         align=PP_ALIGN.CENTER)
    return g


# ================================================================ demo media (cached)
_T = os.path.join(ROOT, "Software", "UTM_PyQt6", "Test data",
                  "8.6.20 - Tensile test to Failure")
_VID = [
    ("demo_pla.mp4", "PLA — S13", "black PLA · fractures at ~4 % strain",
     os.path.join(_T, "Specimen_S13_V2_Spray_Video7", "20260818_111143", "video.avi"), 2.93),
    ("demo_petg.mp4", "PETG — S30", "glossier, more ductile · ~8 % at fracture",
     os.path.join(_T, "Specimen_S30_V3_PETG_Spray_Video9", "20260822_180252", "video.avi"), 4.50),
    ("demo_tpu.mp4", "TPU — S35", "rubber-like · leaves the frame before it fractures",
     os.path.join(_T, "Specimen_S35_V5_TPU_Spray_Video14", "20260824_151610", "video.avi"), 6.64),
]


def ensure_media():
    os.makedirs(MEDIA, exist_ok=True)
    for name, _t, _s, src, speed in _VID:
        dst = os.path.join(MEDIA, name)
        post = dst[:-4] + "_poster.png"
        if not os.path.exists(dst):
            print("  encoding", name)
            subprocess.run(
                ["ffmpeg", "-y", "-v", "error", "-i", src,
                 "-vf", "transpose=1,scale=1280:-2,setpts=PTS/%.2f,fps=30" % speed,
                 "-an", "-c:v", "libx264", "-preset", "veryfast", "-crf", "26",
                 "-movflags", "+faststart", dst], check=True)
        if not os.path.exists(post):
            subprocess.run(["ffmpeg", "-y", "-v", "error", "-ss", "5", "-i", dst,
                            "-frames:v", "1", post], check=True)


# =========================================================================== 01 · title
s = prs.slides.add_slide(BLANK)
base(s, rail=False)
rect(s, 0.30, 0, 0.52, 1.00, INK)
rect(s, 0.30, 1.14, 0.52, EH - 1.14, SAGE)
ju(s, x=1.30, y=0.40, hgt=0.52)
squares(s)
text(s, 1.28, 2.02, 10.6, 2.3,
     [[("PPD-UTM", {})]], fs=66, font=DISP, colour=INK)
text(s, 1.32, 3.22, 10.6, 0.6,
     "A camera-instrumented desktop tensile tester — built, validated against a "
     "commercial lab, and documented.",
     fs=15.5, font=BODY, colour=SAGE)
tricolour(s, 1.32, 3.94, w=4.3, h=0.075)
text(s, 1.32, 5.55, 7.4, 1.5, [
    [("Adithya Sivakumar", {"font": SEMI, "fs": 14})],
    [("adithya.sivakumar@ju.se", {"fs": 11, "colour": SAGE})],
    [("", {"fs": 5})],
    [("Supervisors:  [ supervisor 1 ]  ·  [ supervisor 2 ]", {"fs": 11, "colour": SAGE})],
], spacing=1.15)
vt = text(s, 10.6, 3.3, 2.6, 0.34, "SEPTEMBER 2026", fs=9.5, font=SEMI, colour=INK,
          align=PP_ALIGN.CENTER, wrap=False)
vt.rotation = 270
vt.left, vt.top = Inches(11.85), Inches(3.42)
text(s, 8.6, 6.62, 4.3, 0.4, "Prototype Product Development · Jönköping University",
     fs=10, colour=SAGE, align=PP_ALIGN.RIGHT)
pageno(s)

# ========================================================================== 02 · agenda
s = content_slide("agenda", "What the next 30 minutes cover")
_AG = [
    ("01", "DIC, in one picture", "what the rig measures, and why no sensor touches the part"),
    ("02", "What was built", "hardware, software — and 19 things it does on its own"),
    ("03", "Live demo", "a specimen pulled to fracture, live on the rig"),
    ("04", "Validation", "the same material through a commercial lab — and how close we land"),
    ("05", "AI in the loop", "how this was built by one engineer and a coding agent"),
    ("06", "What comes next", "future work, and one number for fun"),
]
for i, (n, t, d) in enumerate(_AG):
    r, c = divmod(i, 2)
    x, y = 0.85 + c * 6.15, 2.20 + r * 1.45
    card(s, x, y, 5.75, 1.18)
    text(s, x + 0.22, y + 0.16, 0.95, 0.8, n, fs=30, font=DISP, colour=TAN)
    text(s, x + 1.05, y + 0.18, 4.55, 0.4, t, fs=14.5, font=SEMI, colour=INK)
    text(s, x + 1.05, y + 0.60, 4.55, 0.5, d, fs=10.5, colour=SAGE)
text(s, 0.85, 6.72, 11.6, 0.3,
     "Appendix afterwards: Otsu thresholding · PLA vs PETG vs TPU · rig-to-rig noise.",
     fs=10.5, colour=SAGE)

# =========================================================================== 03 · DIC?
s = content_slide("01 · background", "DIC — the camera is the strain gauge", side="BACKGROUND")
pic(s, os.path.join(FIGS, "ppd_dic_explain.png"), 0.85, 2.05, 11.6, 3.05)
_c = [("No contact", "nothing touches the specimen — no clip-on gauge to knock off at fracture"),
      ("No calibration", "strain is a ratio of two pixel distances — mm per pixel cancels out"),
      ("Every frame", "35 measurements per second, live on screen while the test runs")]
for i, (t, d) in enumerate(_c):
    x = 0.85 + i * 3.95
    card(s, x, 5.35, 3.70, 1.45)
    rect(s, x, 5.35, 0.07, 1.45, TAN)
    text(s, x + 0.25, 5.52, 3.3, 0.35, t, fs=13.5, font=SEMI, colour=INK)
    text(s, x + 0.25, 5.90, 3.3, 0.85, d, fs=10.5, colour=INK, spacing=1.05)
text(s, 0.85, 6.95, 11.6, 0.3,
     "DIC = Digital Image Correlation — here in its simplest two-marker form.",
     fs=10, colour=SAGE)

# ==================================================================== 04 · what was built
s = content_slide("02 · the build", "One rig, one application — built together", side="THE BUILD")
pic(s, os.path.join(FIGS, "UTM rig_12-08-26.jpg"), 0.85, 2.00, 3.05, 4.35)
text(s, 0.85, 6.42, 3.05, 0.3, "the PPD-UTM", fs=10, colour=SAGE, align=PP_ALIGN.CENTER)

text(s, 4.25, 2.02, 4.2, 0.34, "HARDWARE", fs=11, font=SEMI, colour=TAN)
bullets(s, 4.25, 2.38, 4.15, 2.0, [
    "twin ball-screw frame, 4.5 kN load cell",
    "industrial camera + LED panel — the DIC eye",
    "printed grips, emergency stop, enclosure",
], fs=11, gap=5)
text(s, 4.25, 4.02, 4.2, 0.34, "SOFTWARE", fs=11, font=SEMI, colour=TAN)
bullets(s, 4.25, 4.38, 4.15, 2.3, [
    "one PyQt6 application runs the whole test",
    "live DIC strain at 35 fps, six closed-loop protocols",
    "every run leaves a CSV + video + PDF report",
    "post-processor re-measures any recorded video",
], fs=11, gap=5)

_stats = [("28", "specimens pulled"), ("3", "materials"), ("19", "smart features"),
          ("~44 k", "lines of Python"), ("204", "evidence pages")]
for i, (n, d) in enumerate(_stats):
    y = 2.02 + i * 0.92
    card(s, 8.85, y, 3.6, 0.80)
    text(s, 9.05, y + 0.10, 1.5, 0.6, n, fs=22, font=DISP, colour=INK)
    text(s, 10.45, y + 0.24, 1.9, 0.5, d, fs=10.5, colour=SAGE)
text(s, 8.85, 6.72, 3.6, 0.3, "numbers as of Sep 2026", fs=9, colour=SAGE)

# ============================================================= 05a · the nineteen features
_SF = [
    (1, "DIC health HUD"), (2, "Prepare specimen"), (3, "Settings / recipes"),
    (4, "1-click PDF report"), (5, "Auto-stop at fracture"), (6, "Strain-rate control"),
    (7, "Stall guard"), (8, "Release load"), (9, "6 test protocols"),
    (10, "Auto-preload"), (11, "Auto-metadata link"), (12, "DIC auto-calibrate"),
    (13, "Guided wizard"), (14, "Test registry"), (15, "Dead-DIC guard"),
    (16, "Live Px₀ overlay"), (17, "Video + stills capture"), (18, "DIC post-processing"),
    (19, "Noise capture"),
]
_HERO = {
    5: "sees the load collapse and halts the motor itself — no hand on the stop button",
    6: "closes the loop on measured strain rate, not motor speed — constant dε/dt",
    18: "re-measures ANY recorded video — it is how the commercial lab's test got audited",
}


def sf_grid(s, heroes=False):
    cols, x0, y0 = 4, 0.85, 2.10
    w, h, gx, gy = 2.83, 0.82, 0.12, 0.14
    for i, (n, name) in enumerate(_SF):
        r, c = divmod(i, cols)
        x, y = x0 + c * (w + gx), y0 + r * (h + gy)
        hero = heroes and n in _HERO
        if heroes and not hero:
            t = card(s, x, y, w, h, fill=CREAM, line=FAINT)
            text(s, x + 0.14, y + 0.10, w - 0.28, 0.3, "SF%d" % n, fs=8.5, font=SEMI,
                 colour=FAINT)
            text(s, x + 0.14, y + 0.32, w - 0.28, 0.42, name, fs=10.5, colour=FAINT)
        elif hero:
            dx, dy = 0.17, 0.13
            t = card(s, x - dx, y - dy, w + 2 * dx, h + 2 * dy, fill=INK, line=TAN, lw=1.4)
            text(s, x - dx + 0.16, y - dy + 0.09, w + 2 * dx - 0.3, 0.3, "SF%d" % n,
                 fs=9, font=SEMI, colour=TAN)
            text(s, x - dx + 0.16, y - dy + 0.30, w + 2 * dx - 0.3, 0.4, name,
                 fs=12.5, font=SEMI, colour=WHITE)
            text(s, x - dx + 0.16, y - dy + 0.60, w + 2 * dx - 0.3, 0.5, _HERO[n],
                 fs=8.2, colour=RGBColor(0xC9, 0xD2, 0xD4), spacing=0.98)
        else:
            t = card(s, x, y, w, h)
            rect(s, x, y, 0.05, h, SAGE)
            text(s, x + 0.16, y + 0.10, w - 0.28, 0.3, "SF%d" % n, fs=8.5, font=SEMI,
                 colour=TAN)
            text(s, x + 0.16, y + 0.32, w - 0.28, 0.42, name, fs=11, colour=INK)
    return t


s = content_slide("02 · the build", "Nineteen things the rig does on its own", side="THE BUILD")
sf_grid(s, heroes=False)
text(s, 0.85, 7.00, 11.6, 0.3,
     "each of these has its own evidence page in the project deck", fs=10, colour=SAGE)

# ============================================================= 05b · the three that matter
s = content_slide("02 · the build", "…and the three worth two minutes each", side="THE BUILD")
sf_grid(s, heroes=True)
text(s, 0.85, 7.00, 11.6, 0.3,
     "SF5 protects the specimen · SF6 makes tests comparable · SF18 made the validation possible",
     fs=10.5, font=SEMI, colour=INK)

# ======================================================================= 06 · live demo
ensure_media()
s = content_slide("03 · live demo", "Three materials, three behaviours — then the real thing",
                  side="LIVE DEMO")
for i, (name, label, sub, _src, _sp) in enumerate(_VID):
    y = 1.98 + i * 1.645
    vw, vh = 9.35, 1.54
    card(s, 0.85, y, vw, vh)
    s.shapes.add_movie(os.path.join(MEDIA, name), Inches(0.94), Inches(y + 0.045),
                       Inches(vw - 0.18), Inches(vh - 0.09),
                       poster_frame_image=os.path.join(MEDIA, name[:-4] + "_poster.png"),
                       mime_type="video/mp4")
    chip(s, 0.85, y - 0.005, label, fill=INK, fs=9)
    text(s, 10.35, y + 0.24, 2.6, 0.4, label.split(" — ")[0], fs=15, font=SEMI, colour=INK)
    text(s, 10.35, y + 0.58, 2.6, 0.9, sub, fs=9.5, colour=SAGE, spacing=1.0)
text(s, 0.85, 6.97, 11.6, 0.3,
     "15-second timelapses of real pulls on this rig (click to play) — then we switch to the "
     "live software and break one for real.", fs=10, colour=SAGE)

# ===================================================================== 07 · validation setup
s = content_slide("04 · validation", "The same material, through a commercial lab",
                  side="VALIDATION")
ghost(s, 0.85, 2.05, 5.6, 4.35, "MOT LAB PHOTO", "drops in here — machine + XT-205 in view")
text(s, 6.85, 2.05, 5.7, 0.34, "THE REFERENCE", fs=11, font=SEMI, colour=TAN)
bullets(s, 6.85, 2.42, 5.55, 1.9, [
    ("MOT materials lab — commercial tensile frame", {"bold": True}),
    "XT-205 video extensometer: a camera that tracks two painted marks — the same principle "
    "as our DIC, in certified form",
    "two specimens from our own print batch: XT205-S1 and XT205-S2",
], fs=11, gap=6)
text(s, 6.85, 4.35, 5.7, 0.34, "TEST CONDITIONS · XT205-S2", fs=11, font=SEMI, colour=TAN)
for i, t in enumerate(("45 mm gauge", "2 mm/min", "300 N preload")):
    chip(s, 6.85 + i * 1.75, 4.72, t, fill=SAGE, fs=10)
bullets(s, 6.85, 5.30, 5.55, 1.4, [
    "same stress-strain quantities, computed two independent ways",
    "and their raw video, re-measured by OUR post-processor (SF18)",
], fs=11, gap=6)

# ===================================================================== 08 · validation results
s = content_slide("04 · validation", "Two rigs, one material — the curves land together",
                  side="VALIDATION")
pic(s, os.path.join(FIGS, "mot2cmp_pair.png"), 0.85, 2.02, 7.45, 4.90)
_tbl = [("UTS", "48.34 / 44.62 MPa", "−7.7 %"),
        ("Yield σy", "39.23 / 39.61 MPa", "+1.0 %"),
        ("Modulus E", "3.45 / 3.23 GPa", "−6.4 %"),
        ("Fracture εf", "5.25 / 5.21 %", "−0.8 %")]
text(s, 8.55, 2.02, 3.9, 0.32, "S34 (PPD-UTM)  /  XT205-S2 (MOT)", fs=10.5, font=SEMI,
     colour=TAN)
for i, (q, v, o) in enumerate(_tbl):
    y = 2.40 + i * 0.62
    card(s, 8.55, y, 3.9, 0.53)
    text(s, 8.72, y + 0.09, 1.35, 0.35, q, fs=11, font=SEMI, colour=INK)
    text(s, 9.95, y + 0.11, 1.75, 0.35, v, fs=9.5, colour=SAGE)
    text(s, 11.65, y + 0.08, 0.72, 0.35, o, fs=11, font=SEMI,
         colour=INK if abs(float(o.strip("%+− ").replace("−", "-"))) < 5 else TAN)
bullets(s, 8.55, 5.15, 3.95, 1.8, [
    "strain agrees to 1 % — the DIC is doing its job",
    "stress gap is real but small — different specimens, different day, 3× pull speed",
    "shape of the whole curve matches to fracture",
], fs=10, gap=5)

# ===================================================================== 09 · can it be trusted
s = content_slide("04 · validation", "Same video, two calculations — one line",
                  side="VALIDATION")
pic(s, os.path.join(FIGS, "mot2pp_agree.png"), 0.85, 2.02, 8.15, 4.90)
_big = [("0.99928", "scale — ours vs theirs", "0.072 % from perfect"),
        ("−55 µε", "constant offset", "= 0.07 px on a 1255 px gauge"),
        ("0.9999964", "R², 3 153 matched instants", "one video, two independent codes")]
for i, (n, t, d) in enumerate(_big):
    y = 2.02 + i * 1.68
    card(s, 9.25, y, 3.22, 1.52)
    rect(s, 9.25, y, 0.07, 1.52, TAN)
    text(s, 9.48, y + 0.13, 2.9, 0.55, n, fs=24, font=DISP, colour=INK)
    text(s, 9.48, y + 0.72, 2.85, 0.35, t, fs=10.5, font=SEMI, colour=INK)
    text(s, 9.48, y + 1.02, 2.85, 0.45, d, fs=9, colour=SAGE)
text(s, 0.85, 7.00, 11.6, 0.3,
     "The MOT lab's own fracture video, fed through our post-processor: the two strain "
     "calculations are the same measurement.", fs=10.5, font=SEMI, colour=INK)

# ===================================================================== 10 · why differences
s = content_slide("04 · validation", "Where the remaining gap comes from", side="VALIDATION")
_why = [
    ("Pull speed", "measured",
     ["their pull: 2 mm/min · ours: ≈3× faster",
      "but at the GAUGE the strain-rate ratio is only 1.19×",
      "→ predicts under 1 % of the stress gap — a real effect, not the story"]),
    ("Grips & seating", "measured",
     ["only 21–45 % of crosshead travel reaches the gauge",
      "our four runs cluster at 21 % and 33 % — a joint that either seats or doesn't",
      "→ affects travel-based numbers, never the DIC strain"]),
    ("Specimen & moisture", "study pending",
     ["different print batches, weeks apart, unconditioned storage",
      "PLA is hygroscopic — literature review queued",
      "→ explanation lands here once the study is done"]),
]
for i, (t, tag, rows) in enumerate(_why):
    x = 0.85 + i * 3.98
    card(s, x, 2.10, 3.72, 3.95)
    rect(s, x, 2.10, 3.72, 0.09, (INK, SAGE, TAN)[i])
    text(s, x + 0.24, 2.32, 3.3, 0.4, t, fs=15, font=SEMI, colour=INK)
    chip(s, x + 0.24, 2.76, tag, fill=(INK if tag == "measured" else TAN), fs=8.5)
    bullets(s, x + 0.24, 3.28, 3.30, 3.1, rows, fs=10, gap=7)
text(s, 0.85, 6.45, 11.6, 0.35,
     "The honest ranking: seating and specimen dominate; speed is measured and small.",
     fs=10.5, font=SEMI, colour=INK)

# ===================================================================== 11 · AI in the loop
s = content_slide("05 · process", "Built by one engineer and a coding agent", side="PROCESS")
_steps = [("intent", "“the rig must stop itself at fracture”"),
          ("prompt", "precise ask + constraints + how to verify"),
          ("Claude Code", "writes app code, tests, analysis"),
          ("verify on rig", "real pulls; failures go straight back in"),
          ("evidence", "every claim becomes a measured deck page")]
for i, (t, d) in enumerate(_steps):
    y = 2.10 + i * 0.99
    card(s, 0.85, y, 4.55, 0.84)
    rect(s, 0.85, y, 0.07, 0.84, TAN if i % 2 else INK)
    text(s, 1.08, y + 0.10, 1.55, 0.4, t, fs=12.5, font=SEMI, colour=INK)
    text(s, 2.65, y + 0.13, 2.7, 0.6, d, fs=9.5, colour=SAGE, spacing=0.98)
    if i < 4:
        text(s, 0.98, y + 0.76, 0.4, 0.3, "↓", fs=12, font=SEMI, colour=SAGE)

_ai = [("348", "commits · 94 % co-authored with the agent"),
       ("9 weeks", "from first DIC frame to a validated rig — 314 commits"),
       ("~44 k", "lines of Python: app, analysis, this deck's generator"),
       ("204", "evidence pages, regenerated from raw data on every build")]
for i, (n, d) in enumerate(_ai):
    y = 2.10 + i * 0.99
    card(s, 5.70, y, 6.78, 0.84)
    text(s, 5.95, y + 0.14, 1.75, 0.55, n, fs=21, font=DISP, colour=INK)
    text(s, 7.55, y + 0.24, 4.8, 0.5, d, fs=10.5, colour=SAGE)
bullets(s, 5.70, 6.12, 6.8, 1.0, [
    ("the prompt is the engineering: what to build, what would falsify it, how to prove it ran",
     {"bold": True}),
    "the agent measures instead of assuming — it renders every page and looks at it",
], fs=10.5, gap=5)

# ===================================================================== 12 · future work
s = content_slide("06 · outlook", "Future work — named, sized, and how", side="OUTLOOK")
_fw = [
    ("2D-DIC strain field", "full-field strain over the speckled gauge, not just two dots",
     "offline first, in the existing post-processor; needs lower exposure — the white body "
     "currently saturates the sensor"),
    ("Correlation, not binary", "sub-pixel tracking that keeps the grey levels",
     "measured gain: 1.1–2.1× less jitter; add as a post-processing option before touching "
     "the live loop"),
    ("Feed simulation models", "measured strain fields calibrating FE material cards",
     "export the field per frame; compare predicted vs measured maps on the same geometry"),
    ("Poisson & true stress (FW1)", "four markers → lateral strain, ν, Cauchy stress",
     "maths written and self-tested; blocked on optics — needs a gauge-zoomed second camera "
     "or a wider specimen"),
    ("Grip-seating experiment (FW2)", "why gauge share clusters at 21 % vs 33 %",
     "one specimen, re-mounted several times, elastic pulls only — spread survives → seating"),
    ("Rate-matched MOT rerun", "close the last confound in the lab comparison",
     "pull at 5.01 mm/min (rate-matched) and 2.00 mm/min (speed-matched) and re-compare"),
]
for i, (t, what, how) in enumerate(_fw):
    r, c = divmod(i, 3)
    x, y = 0.85 + c * 3.98, 2.05 + r * 2.42
    card(s, x, y, 3.72, 2.24)
    rect(s, x, y, 0.07, 2.24, TAN)
    text(s, x + 0.22, y + 0.13, 3.35, 0.4, t, fs=12.5, font=SEMI, colour=INK)
    text(s, x + 0.22, y + 0.50, 3.35, 0.65, what, fs=9.5, colour=INK, spacing=1.0)
    text(s, x + 0.22, y + 1.24, 0.9, 0.3, "HOW", fs=8, font=SEMI, colour=TAN)
    text(s, x + 0.22, y + 1.44, 3.35, 0.75, how, fs=8.8, colour=SAGE, spacing=1.0)

# ===================================================================== 13 · fun facts
s = content_slide("06 · outlook", "One number for fun — the carbon footprint", side="OUTLOOK")
ghost(s, 0.85, 2.20, 11.6, 4.30, "CARBON-FOOTPRINT NUMBERS LAND HERE",
      "the calculation is queued: filament + rig power + compute for the AI collaboration")
text(s, 0.85, 6.75, 11.6, 0.35,
     "spoiler: the printed specimens are not the biggest line item.", fs=10.5, colour=SAGE)

# ===================================================================== 14 · thanks
s = prs.slides.add_slide(BLANK)
base(s, rail=False)
rect(s, 0.30, 0, 0.52, 1.00, INK)
rect(s, 0.30, 1.14, 0.52, EH - 1.14, SAGE)
ju(s, x=1.30, y=0.40, hgt=0.52)
squares(s)
text(s, 1.28, 2.55, 10.6, 1.6, "Thank you.", fs=60, font=DISP, colour=INK)
tricolour(s, 1.32, 3.85, w=4.3, h=0.075)
text(s, 1.32, 4.35, 9.0, 0.5, "Questions — or come break a specimen with me.",
     fs=15, colour=SAGE)
text(s, 1.32, 6.30, 8.0, 0.8, [
    [("Adithya Sivakumar", {"font": SEMI, "fs": 13})],
    [("adithya.sivakumar@ju.se", {"fs": 11, "colour": SAGE})],
], spacing=1.2)
pageno(s)

# ===================================================================== 15 · appendix toc
s = content_slide("appendix", "For the curious", side="APPENDIX")
_ap = [("A1", "Otsu thresholding", "how the camera decides what is a marker"),
       ("A2", "PLA vs PETG vs TPU", "our curves against what literature expects"),
       ("A3", "Noise, rig against rig", "how quiet each instrument is, measured at rest")]
for i, (n, t, d) in enumerate(_ap):
    y = 2.30 + i * 1.42
    card(s, 0.85, y, 11.6, 1.18)
    text(s, 1.10, y + 0.18, 1.2, 0.8, n, fs=28, font=DISP, colour=TAN)
    text(s, 2.30, y + 0.20, 6.0, 0.45, t, fs=15, font=SEMI, colour=INK)
    text(s, 2.30, y + 0.64, 8.5, 0.45, d, fs=10.5, colour=SAGE)

# ===================================================================== 16 · A1 otsu
s = content_slide("appendix · A1", "Otsu — how the camera picks its threshold", side="APPENDIX")
pic(s, os.path.join(FIGS, "otsu_hist.png"), 0.85, 2.05, 7.05, 4.55)
bullets(s, 8.20, 2.15, 4.25, 4.6, [
    ("WHY it exists", {"font": SEMI, "fs": 11.5}),
    "every frame must be split into “marker” and “not marker” — that takes "
    "one brightness number",
    ("WHAT it does", {"font": SEMI, "fs": 11.5}),
    "tries all 254 cuts, keeps the one giving the two tidiest groups — lands in the histogram's "
    "empty valley",
    ("THE CATCH", {"font": SEMI, "fs": 11.5}),
    "it follows the SCENE, not the markers — so the rig scores it against fixed values "
    "(auto-calibrate) instead of trusting it blindly",
], fs=10.5, gap=6)
text(s, 0.85, 6.90, 11.6, 0.3,
     "Otsu (1979), IEEE Trans. SMC 9(1):62–66 · measured here on 120 frames of specimen S13.",
     fs=9.5, colour=SAGE)

# ===================================================================== 17 · A2 materials
s = content_slide("appendix · A2", "Three materials — measured, against expectations",
                  side="APPENDIX")
pic(s, os.path.join(FIGS, "trio_curves.png"), 0.85, 1.98, 11.6, 2.95)
pic(s, os.path.join(FIGS, "petg_expect.png"), 0.85, 5.00, 9.10, 2.02)
bullets(s, 10.12, 5.06, 2.35, 2.0, [
    "ordering exactly as published",
    "TPU never fractured — it left the frame at 19 % strain",
], fs=9.5, gap=6)
text(s, 0.85, 7.06, 11.6, 0.28,
     "published ranges: filament TDS + printed-part literature, as compiled in the project's "
     "material notes.", fs=9, colour=SAGE)

# ===================================================================== 18 · A3 noise
s = content_slide("appendix · A3", "How quiet is each instrument?", side="APPENDIX")
pic(s, os.path.join(FIGS, "mot2cmp_noise.png"), 0.85, 2.02, 8.15, 4.90)
bullets(s, 9.25, 2.15, 3.25, 4.6, [
    ("both rigs sit in the tens of microstrain", {"font": SEMI}),
    "our floor is set by pixel quantisation: noise ∝ 1/Px₀",
    "predicted 1.78× between our two gauge lengths — measured 1.78×",
    "so the noise is understood, not just small",
], fs=10.5, gap=7)

# ============================================================================ save
os.makedirs(SP, exist_ok=True)
prs.save(OUT)
import zipfile
with zipfile.ZipFile(OUT) as z:
    assert z.testzip() is None
    n = len([x for x in z.namelist() if x.startswith("ppt/slides/slide")
             and x.endswith(".xml")])
print("Saved: %s  (%d slides, zip verified)" % (OUT, n))
